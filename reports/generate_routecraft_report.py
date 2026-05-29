from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "路书匠项目报告.pptx"
PPTX_FALLBACK_PATH = OUT_DIR / "路书匠项目报告_修正版.pptx"
HTML_PATH = OUT_DIR / "路书匠项目报告.html"

FONT = "Microsoft YaHei"
TITLE = RGBColor(29, 29, 31)
TEXT = RGBColor(66, 66, 69)
MUTED = RGBColor(110, 110, 115)
BG = RGBColor(245, 245, 247)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(0, 113, 227)
GREEN = RGBColor(24, 128, 86)
ORANGE = RGBColor(196, 116, 32)
RED = RGBColor(180, 48, 44)
LINE = RGBColor(214, 214, 218)

SLIDE_W = 13.333
SLIDE_H = 7.5


def inch(value: float):
    return Inches(value)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    color=TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float = 1.1,
):
    shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_title(slide, index: int, label: str, title: str, subtitle: str | None = None):
    text_box(slide, 0.7, 0.45, 1.0, 0.25, f"{index:02d}", 12, BLUE, True)
    text_box(slide, 1.12, 0.45, 2.6, 0.25, label, 12, MUTED, True)
    text_box(slide, 0.7, 0.78, 7.8, 0.55, title, 28, TITLE, True)
    if subtitle:
        text_box(slide, 0.72, 1.3, 8.8, 0.35, subtitle, 13, MUTED)


def add_footer(slide, index: int):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0.7), inch(7.05), inch(11.95), inch(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    text_box(slide, 0.72, 7.12, 4.8, 0.18, "路书匠 RouteCraft 项目报告 | 2026-05-26", 8, MUTED)
    text_box(slide, 11.95, 7.12, 0.7, 0.18, f"{index:02d}", 8, MUTED, True, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title=None, body=None, accent=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    if title:
        text_box(slide, x + 0.22, y + 0.18, w - 0.44, 0.28, title, 15, TITLE, True)
    if body:
        text_box(slide, x + 0.22, y + 0.55, w - 0.44, h - 0.68, body, 10, TEXT)
    if accent:
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, inch(x + w - 0.42), inch(y + 0.18), inch(0.15), inch(0.15)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
    return shape


def pill(slide, x, y, w, text, color=BLUE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.34)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(238, 246, 255)
    shape.line.color.rgb = RGBColor(205, 226, 252)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = color
    return shape


def bullet_list(slide, x, y, w, h, items, size=14, color=TEXT, gap=0.34):
    for idx, item in enumerate(items):
        cy = y + idx * gap
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x), inch(cy + 0.08), inch(0.08), inch(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        text_box(slide, x + 0.18, cy, w - 0.18, 0.28, item, size, color)


def metric(slide, x, y, w, value, label, color=BLUE):
    card(slide, x, y, w, 1.0)
    text_box(slide, x + 0.2, y + 0.18, w - 0.4, 0.34, value, 26, color, True)
    text_box(slide, x + 0.2, y + 0.62, w - 0.4, 0.22, label, 9, MUTED, True)


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(1, inch(x1), inch(y1), inch(x2), inch(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text_box(slide, 0.72, 0.55, 2.3, 0.3, "RouteCraft", 14, BLUE, True)
    text_box(slide, 0.72, 1.35, 6.4, 0.75, "路书匠项目报告", 38, TITLE, True)
    text_box(
        slide,
        0.75,
        2.15,
        6.6,
        0.65,
        "轻量 AI 出行规划工具：从出行条件输入，到外部数据增强，再到大模型流式生成可执行路线建议。",
        16,
        TEXT,
    )
    pill(slide, 0.78, 3.0, 1.25, "AI 规划")
    pill(slide, 2.18, 3.0, 1.35, "SSE 流式")
    pill(slide, 3.68, 3.0, 1.15, "高德")
    pill(slide, 4.98, 3.0, 1.15, "天气")
    pill(slide, 6.28, 3.0, 1.35, "实时检索")

    # Abstract route canvas
    canvas = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(7.55), inch(0.82), inch(4.95), inch(5.78)
    )
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = WHITE
    canvas.line.color.rgb = LINE
    for x, y, label, col in [
        (8.0, 1.5, "起点", BLUE),
        (10.05, 2.45, "天气", ORANGE),
        (8.82, 3.55, "景点", GREEN),
        (11.2, 4.72, "终点", BLUE),
    ]:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x), inch(y), inch(0.52), inch(0.52))
        node.fill.solid()
        node.fill.fore_color.rgb = col
        node.line.fill.background()
        text_box(slide, x - 0.1, y + 0.62, 0.72, 0.2, label, 8, MUTED, True, PP_ALIGN.CENTER)
    arrow(slide, 8.55, 1.75, 9.9, 2.55, RGBColor(120, 153, 184))
    arrow(slide, 10.2, 2.98, 9.13, 3.45, RGBColor(120, 153, 184))
    arrow(slide, 9.4, 3.82, 11.0, 4.9, RGBColor(120, 153, 184))
    text_box(slide, 8.0, 5.65, 3.95, 0.45, "规划链路可追溯：输入快照、外部数据快照、流式事件、最终结果均落库。", 12, TEXT)
    add_footer(slide, 1)


def add_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 2, "项目结论", "当前阶段结论", "首版核心链路已经成形，剩余重点是联调收口和真实 LLM Key 修复。")
    metric(slide, 0.78, 1.78, 2.35, "3 端", "后端 API / 用户端 / 管理后台", BLUE)
    metric(slide, 3.42, 1.78, 2.35, "60", "后端 pytest 用例通过", GREEN)
    metric(slide, 6.06, 1.78, 2.35, "14", "MVP 业务表已设计并迁移", ORANGE)
    metric(slide, 8.7, 1.78, 2.35, "P0/P1", "首版接口基本注册完成", BLUE)
    card(
        slide,
        0.78,
        3.1,
        5.25,
        2.25,
        "已具备的核心能力",
        "用户注册登录、游客会话、JWT 鉴权、规划记录、SSE 流式生成、生成输出落库、地图/天气/实时检索适配、管理端记录和 LLM 配置管理。",
        GREEN,
    )
    card(
        slide,
        6.35,
        3.1,
        5.25,
        2.25,
        "待收口的关键风险",
        "真实 LLM 配置中存量 API Key 仍需重新保存为可解密密文；前后端联调、mypy、地图截图兜底和配置审计日志仍需补齐。",
        ORANGE,
    )
    bullet_list(
        slide,
        0.9,
        5.72,
        10.8,
        0.9,
        [
            "项目定位清晰：先验证 AI 出行规划体验，不引入订单、支付、复杂 RBAC。",
            "技术架构与数据快照策略支持历史回放、失败恢复和后台排障。",
        ],
        13,
    )
    add_footer(slide, 2)


def add_product_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 3, "产品定位", "MVP 范围与用户价值", "输入少量约束，输出可直接参考的行程建议。")
    card(slide, 0.78, 1.75, 3.35, 3.85, "用户输入", "必填：起点、目的地、范围、交通方式\n选填：日期、人数、偏好、避免项\n\n目标：降低规划启动门槛，避免复杂下单流程。", BLUE)
    card(slide, 4.5, 1.75, 3.35, 3.85, "系统输出", "需求理解、天气预警、路径点/公共交通规划、高德路线链接和路径图、途径景点、实时信息、时间安排、风险提示、Markdown 与 JSON 摘要。", GREEN)
    card(slide, 8.22, 1.75, 3.35, 3.85, "角色边界", "用户/游客：生成规划、看流式过程、查看历史。\n管理员：用户管理、记录管理、LLM 配置、异常处理。\n\n不做规划师、运营、支付和订单。", ORANGE)
    text_box(slide, 0.9, 6.05, 10.8, 0.36, "一句话定位：用户输入出行条件后，系统通过大模型流式生成包含天气、路线、景点和实时资讯的出行规划。", 15, TITLE, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)


def add_user_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 4, "核心流程", "从输入到结果的生成链路", "SSE 让用户持续看到规划过程，数据库保留完整事件和结果。")
    steps = [
        ("01", "提交条件", "起点、目的地、范围、交通方式"),
        ("02", "创建记录", "生成记录与输入快照落库"),
        ("03", "聚合数据", "高德、天气、实时检索"),
        ("04", "LLM 编排", "拼装上下文并流式调用"),
        ("05", "阶段输出", "token / snapshot / done"),
        ("06", "保存结果", "Markdown、JSON、快照、错误"),
    ]
    x = 0.82
    for i, (no, title, desc) in enumerate(steps):
        card(slide, x + i * 1.95, 2.05, 1.55, 2.35)
        text_box(slide, x + i * 1.95 + 0.18, 2.28, 0.48, 0.25, no, 11, BLUE, True)
        text_box(slide, x + i * 1.95 + 0.18, 2.72, 1.05, 0.36, title, 15, TITLE, True)
        text_box(slide, x + i * 1.95 + 0.18, 3.25, 1.08, 0.66, desc, 9, TEXT)
        if i < len(steps) - 1:
            arrow(slide, x + i * 1.95 + 1.55, 3.2, x + (i + 1) * 1.95 - 0.08, 3.2, RGBColor(146, 160, 174))
    card(
        slide,
        1.05,
        5.15,
        10.7,
        0.85,
        None,
        "流式事件：record_created → stage → token → snapshot → done / error。连接中断后可通过记录详情或续接接口恢复已保存内容。",
        BLUE,
    )
    add_footer(slide, 4)


def add_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 5, "技术架构", "三端协作与外部能力接入", "FastAPI 作为统一后端，用户端与管理端共享认证和记录能力。")
    card(slide, 0.78, 1.55, 2.65, 1.2, "Web 用户端", "Vue 3 / Vite / Pinia / Element Plus\n端口 3003", BLUE)
    card(slide, 0.78, 3.25, 2.65, 1.2, "管理后台", "Vue 3 / Vite / 路由守卫\n端口 3004", ORANGE)
    card(slide, 4.35, 2.05, 3.25, 2.1, "后端 API", "FastAPI / SQLAlchemy 2 / Pydantic v2\n认证、记录、SSE、外部服务编排\n端口 3002", BLUE)
    card(slide, 8.45, 1.25, 3.05, 1.05, "MySQL", "用户、记录、快照、日志、配置", GREEN)
    card(slide, 8.45, 2.62, 3.05, 1.05, "Redis / Cache", "限流与外部 API TTL 缓存", GREEN)
    card(slide, 8.45, 4.0, 3.05, 1.45, "外部服务", "高德 Web 服务\n腾讯天气 / 高德天气\nTavily 实时检索\nOpenAI-compatible LLM", ORANGE)
    arrow(slide, 3.45, 2.15, 4.28, 2.7)
    arrow(slide, 3.45, 3.85, 4.28, 3.35)
    arrow(slide, 7.68, 2.7, 8.38, 1.8, RGBColor(120, 153, 184))
    arrow(slide, 7.68, 3.0, 8.38, 3.08, RGBColor(120, 153, 184))
    arrow(slide, 7.68, 3.35, 8.38, 4.55, RGBColor(120, 153, 184))
    bullet_list(
        slide,
        0.9,
        6.0,
        10.7,
        0.55,
        ["统一响应结构、CORS、请求日志和链路 ID 已接入；外部接口无 Key 或配置 Mock 时运行期明确失败。"],
        12,
    )
    add_footer(slide, 5)


def add_ai_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 6, "AI 编排", "规划生成的 8 个阶段", "阶段化输出让用户可感知进度，也方便后端落库和排障。")
    stages = [
        ("需求理解", BLUE),
        ("天气预警", ORANGE),
        ("地图 POI", GREEN),
        ("路线规划", BLUE),
        ("路径图", GREEN),
        ("途径景点", GREEN),
        ("实时信息", ORANGE),
        ("汇总建议", BLUE),
    ]
    for i, (name, color) in enumerate(stages):
        row = i // 4
        col = i % 4
        x = 0.78 + col * 2.95
        y = 1.75 + row * 1.72
        card(slide, x, y, 2.35, 1.15, f"{i + 1}. {name}", "输入上下文、生成阶段内容、保存阶段事件。", color)
    card(
        slide,
        0.78,
        5.28,
        5.25,
        1.12,
        "LLM 输出",
        "最终 Markdown + JSON 摘要，承载天气、路线、景点、实时资讯和风险。",
        BLUE,
    )
    card(
        slide,
        6.35,
        5.28,
        5.25,
        1.12,
        "可追溯上下文",
        "TripPlanningContext 汇总 weather、route、transport、map_export、attractions、realtime、risks。",
        GREEN,
    )
    add_footer(slide, 6)


def add_web_app(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 7, "用户端", "单页规划工作台", "首屏直接进入规划，不做营销页。")
    card(slide, 0.78, 1.55, 3.15, 4.55, "规划输入", "步骤式表单：去哪儿、怎么去、什么时候、偏好与避开、确认信息。\n\n支持省市县选择、交通方式、人数步进器、多选偏好标签。", BLUE)
    card(slide, 4.23, 1.55, 3.15, 4.55, "流式输出", "状态栏、阶段进度、停止生成、复制结果。\n\ntoken 实时追加，snapshot 更新天气、路线、地图、景点、实时信息模块。", GREEN)
    card(slide, 7.68, 1.55, 3.15, 4.55, "历史与详情", "生成完成后保留摘要和结构化结果。\n\n历史记录支持进入详情、重新生成、失败重试和流式过程续接。", ORANGE)
    text_box(slide, 0.95, 6.32, 10.6, 0.28, "视觉基调：浅灰页面背景、白色内容面、近黑主文字、蓝色行动色，强调克制、清爽和可读性。", 12, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide, 7)


def add_admin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 8, "管理后台", "运营与排障工作台", "后台聚焦用户、生成记录和 LLM 配置，不引入复杂权限模型。")
    metric(slide, 0.78, 1.62, 2.7, "用户", "列表、搜索、状态筛选、启用/禁用", BLUE)
    metric(slide, 3.78, 1.62, 2.7, "记录", "筛选、详情、错误、重试、删除", GREEN)
    metric(slide, 6.78, 1.62, 2.7, "LLM", "配置、密钥掩码、启用/停用、连接测试", ORANGE)
    card(
        slide,
        0.78,
        3.25,
        5.25,
        2.3,
        "控制台概览",
        "统计总用户数、生成记录数、已启用 LLM 配置数，便于快速判断当前系统资源规模。",
        BLUE,
    )
    card(
        slide,
        6.35,
        3.25,
        5.25,
        2.3,
        "安全边界",
        "管理端必须校验 role=admin；API Key 只展示掩码，日志和响应不得返回完整 Token、密码或密钥。",
        RED,
    )
    add_footer(slide, 8)


def add_data_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 9, "数据设计", "以快照保证历史可复现", "用户输入、外部数据、流式事件和最终结果都独立保存。")
    cols = [
        ("账号与会话", ["users", "login_sessions"], BLUE),
        ("生成记录", ["generation_records", "generation_inputs", "generation_outputs", "generation_stream_events", "generation_errors"], GREEN),
        ("外部快照", ["route_snapshots", "route_map_exports", "weather_snapshots", "news_snapshots", "external_api_cache"], ORANGE),
        ("LLM 与审计", ["llm_configs", "llm_call_logs", "config_audit_logs"], RED),
    ]
    for i, (title, tables, color) in enumerate(cols):
        x = 0.78 + i * 2.95
        card(slide, x, 1.65, 2.38, 4.35, title, None, color)
        y = 2.3
        for table in tables:
            pill(slide, x + 0.22, y, 1.9, table, color)
            y += 0.48
    text_box(slide, 0.94, 6.35, 10.5, 0.3, "索引围绕用户历史、管理端筛选、记录详情、流式恢复和 LLM 成本统计设计。", 12, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide, 9)


def add_status_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 10, "进度与风险", "当前状态与下一步", "后端自测已通过，下一阶段应集中在真实联调和质量门禁。")
    card(slide, 0.78, 1.58, 3.35, 4.72, "已完成", "FastAPI 工程与迁移\n认证、JWT、游客会话\n生成记录、SSE、取消/重试\n高德、天气、Tavily 适配\nAI 规划编排与输出落库\n用户端和管理端首版页面", GREEN)
    card(slide, 4.5, 1.58, 3.35, 4.72, "阻塞与风险", "存量 LLM Key 旧哈希不可解密\n前端联调仍需闭环验证\nmypy 未配置\n路径图截图兜底未实现\nLLM 配置审计日志待补\n供应商错误码需细分", ORANGE)
    card(slide, 8.22, 1.58, 3.35, 4.72, "建议下一步", "重新保存默认 LLM API Key\n跑通真实生成端到端流程\n完成用户端历史续接与重试联调\n补 mypy 或明确暂缓\n补截图兜底与审计日志\n整理演示用种子数据", BLUE)
    add_footer(slide, 10)


def add_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, 11, "交付路线", "MVP 收口计划", "建议用三步完成从可用到可演示的闭环。")
    lanes = [
        ("第 1 步：联调可用", "修复 LLM Key，完成真实生成、详情、历史、重试、管理端记录筛选。", BLUE),
        ("第 2 步：稳定可测", "补充 mypy / 前端 lint 与构建检查，细化外部服务失败提示和供应商错误码。", GREEN),
        ("第 3 步：演示可讲", "准备演示数据、典型路线案例、后台异常处理案例和用户端生成录屏。", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(lanes):
        y = 1.75 + i * 1.55
        card(slide, 1.08, y, 10.4, 1.08, title, desc, color)
        if i < len(lanes) - 1:
            arrow(slide, 6.25, y + 1.08, 6.25, y + 1.42, RGBColor(146, 160, 174))
    text_box(slide, 1.18, 6.4, 10.2, 0.28, "MVP 原则保持不变：先验证规划体验、生成速度和结果可用性，再扩展复杂业务闭环。", 13, TITLE, True, PP_ALIGN.CENTER)
    add_footer(slide, 11)


SLIDES_HTML = [
    ("01 项目封面", "路书匠项目报告", "轻量 AI 出行规划工具：从出行条件输入，到外部数据增强，再到大模型流式生成可执行路线建议。", ["AI 规划", "SSE 流式", "高德", "天气", "实时检索"]),
    ("02 项目结论", "当前阶段结论", "首版核心链路已经成形，剩余重点是联调收口和真实 LLM Key 修复。", ["3 端：后端 API / 用户端 / 管理后台", "60 个后端 pytest 用例通过", "14 张 MVP 业务表已设计并迁移", "P0/P1 首版接口基本注册完成"]),
    ("03 产品定位", "MVP 范围与用户价值", "输入少量约束，输出可直接参考的行程建议。", ["用户输入：起点、目的地、范围、交通方式、日期、人数、偏好、避免项", "系统输出：天气、路线、地图、景点、实时资讯、风险提示、Markdown 与 JSON 摘要", "角色边界：用户/游客与管理员，不做规划师、运营、支付和订单"]),
    ("04 核心流程", "从输入到结果的生成链路", "SSE 让用户持续看到规划过程，数据库保留完整事件和结果。", ["提交条件", "创建记录", "聚合数据", "LLM 编排", "阶段输出", "保存结果"]),
    ("05 技术架构", "三端协作与外部能力接入", "FastAPI 作为统一后端，用户端与管理端共享认证和记录能力。", ["Web 用户端：Vue 3 / Vite / Pinia / Element Plus", "后端 API：FastAPI / SQLAlchemy 2 / Pydantic v2", "MySQL：用户、记录、快照、日志、配置", "外部服务：高德、天气、Tavily、OpenAI-compatible LLM"]),
    ("06 AI 编排", "规划生成的 8 个阶段", "阶段化输出让用户可感知进度，也方便后端落库和排障。", ["需求理解", "天气预警", "地图 POI", "路线规划", "路径图", "途径景点", "实时信息", "汇总建议"]),
    ("07 用户端", "单页规划工作台", "首屏直接进入规划，不做营销页。", ["规划输入：步骤式表单和交通方式选择", "流式输出：阶段进度、停止生成、复制结果", "历史与详情：重新生成、失败重试、流式续接"]),
    ("08 管理后台", "运营与排障工作台", "后台聚焦用户、生成记录和 LLM 配置，不引入复杂权限模型。", ["用户管理：列表、搜索、状态筛选、启用/禁用", "记录管理：筛选、详情、错误、重试、删除", "LLM 配置：密钥掩码、启用/停用、连接测试"]),
    ("09 数据设计", "以快照保证历史可复现", "用户输入、外部数据、流式事件和最终结果都独立保存。", ["账号与会话", "生成记录", "外部快照", "LLM 与审计"]),
    ("10 进度与风险", "当前状态与下一步", "后端自测已通过，下一阶段应集中在真实联调和质量门禁。", ["已完成：工程、认证、SSE、记录、外部集成、AI 编排、首版页面", "风险：LLM Key、前端联调、mypy、截图兜底、审计日志", "下一步：真实端到端联调、质量门禁、演示数据"]),
    ("11 交付路线", "MVP 收口计划", "建议用三步完成从可用到可演示的闭环。", ["第 1 步：联调可用", "第 2 步：稳定可测", "第 3 步：演示可讲"]),
]


def build_html():
    slides = []
    for idx, (label, title, subtitle, items) in enumerate(SLIDES_HTML, start=1):
        item_html = "".join(f"<li>{item}</li>" for item in items)
        slides.append(
            f"""
            <section class="slide" data-screen-label="{label}">
              <div class="topline"><span>{idx:02d}</span><b>{label.split(' ', 1)[1]}</b></div>
              <h1>{title}</h1>
              <p class="subtitle">{subtitle}</p>
              <ul class="grid">{item_html}</ul>
              <footer>路书匠 RouteCraft 项目报告 | 2026-05-26 <span>{idx:02d}</span></footer>
            </section>
            """
        )
    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>路书匠项目报告</title>
  <style>
    :root {{
      --bg: #f5f5f7;
      --surface: #ffffff;
      --text: #1d1d1f;
      --muted: #6e6e73;
      --line: #d6d6da;
      --blue: #0071e3;
      --green: #188056;
      --orange: #c47420;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      min-height: 100vh;
      background: #e8e8ed;
      color: var(--text);
      font-family: "Microsoft YaHei", "PingFang SC", "Noto Sans CJK SC", sans-serif;
      overflow: hidden;
      letter-spacing: 0;
    }}
    .deck-wrap {{
      position: fixed;
      inset: 0;
      display: grid;
      place-items: center;
      padding: 52px 88px 72px;
    }}
    .deck {{
      width: 1920px;
      height: 1080px;
      transform-origin: center center;
      position: relative;
      background: var(--bg);
      box-shadow: 0 28px 70px rgba(15, 23, 42, 0.22);
      overflow: hidden;
    }}
    .slide {{
      position: absolute;
      inset: 0;
      display: none;
      padding: 86px 112px;
      background: var(--bg);
    }}
    .slide.active {{
      display: block;
      animation: fade 180ms ease-out;
    }}
    @keyframes fade {{
      from {{ opacity: 0; transform: translateY(10px); }}
      to {{ opacity: 1; transform: translateY(0); }}
    }}
    .topline {{
      display: flex;
      align-items: center;
      gap: 18px;
      color: var(--muted);
      font-size: 22px;
      font-weight: 700;
    }}
    .topline span {{ color: var(--blue); }}
    h1 {{
      margin: 44px 0 18px;
      width: 1180px;
      font-size: 74px;
      line-height: 1.08;
      text-wrap: pretty;
      letter-spacing: 0;
    }}
    .subtitle {{
      width: 1180px;
      margin: 0;
      color: #424245;
      font-size: 31px;
      line-height: 1.48;
      text-wrap: pretty;
    }}
    .grid {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 24px;
      margin: 96px 0 0;
      padding: 0;
      list-style: none;
      width: 1510px;
    }}
    .grid li {{
      min-height: 122px;
      padding: 28px 32px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--surface);
      color: #2f3337;
      font-size: 27px;
      line-height: 1.42;
      text-wrap: pretty;
    }}
    .grid li::before {{
      content: "";
      display: inline-block;
      width: 11px;
      height: 11px;
      margin-right: 14px;
      border-radius: 50%;
      background: var(--blue);
      vertical-align: middle;
    }}
    footer {{
      position: absolute;
      left: 112px;
      right: 112px;
      bottom: 54px;
      padding-top: 18px;
      border-top: 1px solid var(--line);
      color: var(--muted);
      font-size: 18px;
    }}
    footer span {{ float: right; font-weight: 700; }}
    .nav {{
      position: fixed;
      left: 50%;
      bottom: 18px;
      transform: translateX(-50%);
      display: flex;
      align-items: center;
      gap: 10px;
      z-index: 10;
    }}
    .nav button {{
      min-width: 86px;
      height: 38px;
      border: 1px solid #c9c9cf;
      border-radius: 8px;
      background: #fff;
      color: #1d1d1f;
      font-size: 14px;
      cursor: pointer;
    }}
    .nav button:focus-visible {{
      outline: 2px solid var(--blue);
      outline-offset: 2px;
    }}
    .counter {{
      min-width: 74px;
      color: #424245;
      font-size: 14px;
      text-align: center;
    }}
  </style>
</head>
<body>
  <main class="deck-wrap">
    <div class="deck" id="deck">
      {''.join(slides)}
    </div>
  </main>
  <div class="nav" aria-label="幻灯片导航">
    <button id="prev">上一页</button>
    <span class="counter" id="counter"></span>
    <button id="next">下一页</button>
  </div>
  <script>
    const slides = [...document.querySelectorAll('.slide')];
    const deck = document.getElementById('deck');
    const counter = document.getElementById('counter');
    let index = Number(localStorage.getItem('routecraft-report-slide') || 0);
    index = Math.min(Math.max(index, 0), slides.length - 1);

    function fit() {{
      const padX = 176;
      const padY = 124;
      const scale = Math.min((innerWidth - padX) / 1920, (innerHeight - padY) / 1080);
      deck.style.transform = `scale(${{Math.max(0.18, scale)}})`;
    }}
    function show(nextIndex) {{
      index = Math.min(Math.max(nextIndex, 0), slides.length - 1);
      slides.forEach((slide, i) => slide.classList.toggle('active', i === index));
      counter.textContent = `${{index + 1}} / ${{slides.length}}`;
      localStorage.setItem('routecraft-report-slide', String(index));
    }}
    document.getElementById('prev').addEventListener('click', () => show(index - 1));
    document.getElementById('next').addEventListener('click', () => show(index + 1));
    addEventListener('keydown', (event) => {{
      if (event.key === 'ArrowLeft') show(index - 1);
      if (event.key === 'ArrowRight' || event.key === ' ') {{
        event.preventDefault();
        show(index + 1);
      }}
    }});
    addEventListener('resize', fit);
    fit();
    show(index);
  </script>
</body>
</html>
"""
    HTML_PATH.write_text(html, encoding="utf-8")


def build_pptx(path: Path = PPTX_PATH):
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    add_cover(prs)
    add_executive_summary(prs)
    add_product_scope(prs)
    add_user_flow(prs)
    add_architecture(prs)
    add_ai_pipeline(prs)
    add_web_app(prs)
    add_admin(prs)
    add_data_model(prs)
    add_status_risks(prs)
    add_roadmap(prs)
    prs.save(path)


def main():
    build_html()
    try:
        build_pptx(PPTX_PATH)
        pptx_path = PPTX_PATH
    except PermissionError:
        build_pptx(PPTX_FALLBACK_PATH)
        pptx_path = PPTX_FALLBACK_PATH
    print(pptx_path)
    print(HTML_PATH)


if __name__ == "__main__":
    main()
